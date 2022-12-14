import base64
import glob
import io
import json
import os
import re
from datetime import datetime
from typing import List

import fastapi
import uvicorn as uvicorn
from fastapi import Query
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
from pptx import Presentation
from starlette.middleware.cors import CORSMiddleware
from starlette.requests import Request
from starlette.responses import RedirectResponse, Response, HTMLResponse
from starlette.staticfiles import StaticFiles
from starlette.templating import Jinja2Templates

from database.db_handler import DatabaseHandler, Users, Presentations, Slides, Tags, Links
from utils import utils
import aspose.slides as slides

app = fastapi.FastAPI()

app.mount('/views', StaticFiles(directory="views"), name="views")

origins = [
    "http://localhost",
    "http://localhost:8080",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

db_handler = DatabaseHandler()
current_user_id = None

flow = Flow.from_client_secrets_file(
        'auth/oauth2.keys.json',
        scopes=['openid',
                'https://www.googleapis.com/auth/userinfo.email',
                'https://www.googleapis.com/auth/userinfo.profile',
                'https://www.googleapis.com/auth/drive',
                'https://www.googleapis.com/auth/presentations'],
        redirect_uri="http://localhost:8000/auth/sessions"
    )





@app.get('/')
def root():
    with open(os.path.join('views/index.html'), encoding='utf-8') as fh:
        data = fh.read()
    return HTMLResponse(content=data)


@app.get('/auth/login')
def login():
    auth_uri = flow.authorization_url()[0]
    return RedirectResponse(url=auth_uri)


@app.get('/auth/sessions')
def sessions(code):
    global current_user_id
    flow.fetch_token(code=code)
    credentials = flow.credentials

    user_info_service = build('oauth2', 'v2', credentials=credentials)
    user_info = user_info_service.userinfo().get().execute()
    session = db_handler.Session()
    user = session.query(Users).get(user_info['id'])
    current_user_id = user_info['id']
    if user is None:
        session.add(Users(
            id=user_info.get('id'),
            name=user_info.get('name'),
            refresh_token=credentials.refresh_token,
        ))
        session.commit()


@app.get('/presentations/folders')
def pres_folders():
    credentials = flow.credentials
    drive = build('drive', 'v3', credentials=credentials)
    q = "mimeType='application/vnd.google-apps.folder' and 'root' in parents " \
        "and trashed = false"
    response = drive.files().list(q=q).execute()
    return response.get('files')


@app.get('/presentations/list')
def pres_list(folder_ids: List[str] = Query(...)):
    credentials = flow.credentials
    drive = build('drive', 'v3', credentials=credentials)
    q = "mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation' " \
        "and trashed = false"
    fields = 'files(id, name, modifiedTime)'
    files = []
    for folder_id in folder_ids:
        folder_query = q + f" and '{folder_id}' in parents"
        response = drive.files().list(q=folder_query, fields=fields).execute()
        folder_files = {
            'folder_id': folder_id,
            'presentations': response['files']
        }
        update_presentations(folder_files['folder_id'], folder_files['presentations'])
        files.append(folder_files)

    return files


@app.get('/presentations/slides')
def get_pres_slides(pres_id):
    session = db_handler.Session()
    slides = session.query(Slides).filter(Slides.pres_id == pres_id).all()
    slides = list(map(lambda slide:
                      {
                          'id': slide.id,
                          'pres_id': slide.pres_id,
                          'index': slide.index,
                          'thumbnail': base64.b64encode(slide.thumbnail).decode()
                      }, slides))
    return slides


@app.get('/presentations/slides-by-query')
def get_slides_by_query(query):
    tag_names = re.findall('[a-z]+[a-z0-9]*', query)
    tag_names = list(filter(lambda x: x not in ['and', 'or', 'not'], tag_names))
    session = db_handler.Session()
    links_and_tags = session.query(Links, Tags)\
        .filter(Tags.name.in_(tag_names), Tags.owner_id == current_user_id)\
        .filter(Links.tag_id == Tags.id)\
        .all()
    slides_tags_values = dict()
    for [link, tag] in links_and_tags:
        if slides_tags_values.get(link.slide_id) is None:
            slides_tags_values[link.slide_id] = {key: False for key in tag_names}
        if link.value:
            slides_tags_values[link.slide_id][tag.name] = link.value
        else:
            slides_tags_values[link.slide_id][tag.name] = True
    filtered_slides = []
    for slide_id in slides_tags_values.keys():
        eval_query = query
        for tag in slides_tags_values[slide_id].keys():
            eval_query = eval_query.replace(tag, str(slides_tags_values[slide_id][tag]))
        if eval(eval_query):
            filtered_slides.append(slide_id)
    result = session.query(Slides).filter(Slides.id.in_(filtered_slides)).all()
    result = list(map(lambda x: {
        'id': x.id,
        'pres_id': x.pres_id,
        'index': x.index,
        'thumbnail': base64.b64encode(x.thumbnail).decode()}, result))
    return result


@app.post('/presentations/build')
def build_presentation(new_pres_info: dict):
    clear_temp()

    name = new_pres_info.get('name')
    build_from = new_pres_info.get('body')
    for pres in build_from:
        download_pres(pres.get('id'), f'{pres.get("id")}.pptx')
    with slides.Presentation() as presentation:
        presentation.slide_size.set_size(slides.SlideSizeType.ON_SCREEN_16X9, slides.SlideSizeScaleType.DO_NOT_SCALE)
        presentation.slides.remove_at(0)
        for pres_from in build_from:
            with slides.Presentation(f"presentations/{pres_from.get('id')}.pptx") as presentation_from:
                for slide_index in pres_from.get('slides'):
                    mgr = presentation_from.slides[slide_index].notes_slide_manager
                    mgr.remove_notes_slide()
                    presentation.slides.add_clone(presentation_from.slides[slide_index])
        presentation.save(f"presentations/{name}.pptx", slides.export.SaveFormat.PPTX)
    utils.clear_watermark(f"presentations/{name}.pptx")

    credentials = flow.credentials
    drive = build('drive', 'v3', credentials=credentials)
    file_metadata = {
        'name': name
    }
    media = MediaFileUpload(f"presentations/{name}.pptx",
                            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    file = drive.files().create(
        body=file_metadata,
        media_body=media,
        fields='id, name, modifiedTime, parents').execute()

    update_presentations(file.get('parents')[0], [file], download=False)

    session = db_handler.Session()
    migrate_links = {}
    index = 0
    for pres in build_from:
        pres_id = pres.get('id')
        slides_indices = pres.get('slides')
        slides_from = session.query(Slides).filter(Slides.pres_id == pres_id)
        for i in slides_indices:
            slide_from = slides_from.filter(Slides.index == i).first()
            links = session.query(Links).filter(Links.slide_id == slide_from.id).all()
            migrate_links[index] = links
            index += 1

    slides_to = session.query(Slides).filter(Slides.pres_id == file.get('id'))
    for index in migrate_links.keys():
        slide_to = slides_to.filter(Slides.index == index).first()
        for link in migrate_links[index]:
            session.add(Links(
                slide_id=slide_to.id,
                tag_id=link.tag_id,
                value=link.value
            ))
    session.commit()

    return file.get('id')


@app.post('/links/create')
def create_link(slide_id, tag_name, value=None):
    session = db_handler.Session()
    tag = session.query(Tags).filter(Tags.name == tag_name and Tags.owner_id == current_user_id).first()
    if tag is None:
        session.add(Tags(
            name=tag_name,
            owner_id=current_user_id
        ))
        session.commit()
    tag = session.query(Tags).filter(Tags.name == tag_name, Tags.owner_id == current_user_id).first()
    link = session.query(Links).filter(Links.slide_id == slide_id, Links.tag_id == tag.id).first()
    if link is None:
        session.add(Links(
            slide_id=slide_id,
            tag_id=tag.id,
            value=value
        ))
    else:
        link.value = value
    session.commit()


@app.post('/links/remove')
def remove_link(slide_id, tag_name):
    session = db_handler.Session()
    tag = session.query(Tags).filter(Tags.name == tag_name and Tags.owner_id == current_user_id).first()
    if tag is None:
        return "Tag does not exists"
    link = session.query(Links).filter(Links.slide_id == slide_id and Links.tag_id == tag.id).first()
    if link is None:
        return "Link does not exists"
    session.delete(link)
    session.commit()


@app.get('/links/slide-tags')
def get_slide_tags(slide_id):
    session = db_handler.Session()
    links = session.query(Links).filter(Links.slide_id == slide_id).all()
    tags = []
    for link in links:
        tag = session.query(Tags).get(link.tag_id)
        tags.append({'id': tag.id, 'name': tag.name, 'owner_id': tag.owner_id, 'value': link.value})
    return tags


def download_pres(pres_id, name='temp'):
    try:
        credentials = flow.credentials
        drive = build('drive', 'v3', credentials=credentials)
        request = drive.files().get_media(fileId=pres_id)
        file = io.BytesIO()
        downloader = MediaIoBaseDownload(file, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()

    except HttpError as error:
        print(F'An error occurred: {error}')
        file = None

    with open(f"presentations/{name}", "wb") as f:
        f.write(file.getvalue())


def upload_pres(pres_id, name):
    credentials = flow.credentials
    drive = build('drive', 'v2', credentials=credentials)

    file = drive.files().get(fileId=pres_id).execute()

    media_body = MediaFileUpload(
        f'presentations/{name}',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        resumable=True
    )

    file = drive.files().update(
        fileId=pres_id,
        body=file,
        media_body=media_body
    ).execute()

    session = db_handler.Session()
    session.query(Presentations)\
        .filter(Presentations.id == file['id'])\
        .update({'modified_time': file['modifiedDate']})
    session.commit()


def update_presentations(folder_id, files, download=True):
    if download: clear_temp()
    session = db_handler.Session()
    for pres in files:
        found_in_db = session.query(Presentations).get(pres.get('id'))
        if found_in_db is None:
            session.add(Presentations(
                id=pres.get('id'),
                name=pres.get('name'),
                owner_id=current_user_id,
                folder_id=folder_id,
                modified_time=pres.get('modifiedTime')
            ))
            session.commit()
            update_slides(pres.get('id'), pres.get('name'))
        elif found_in_db.modified_time < datetime.strptime(pres.get('modifiedTime'), '%Y-%m-%dT%H:%M:%S.%fZ'):
            found_in_db.name = pres.get('name')
            found_in_db.modified_time = pres.get('modifiedTime')
            session.commit()
            update_slides(pres.get('id'), pres.get('name'), download)
    session.query(Presentations).filter(
        Presentations.folder_id == folder_id,
        ~Presentations.id.in_([x['id'] for x in files])).delete()
    session.commit()


def update_slides(pres_id, name, download=True):
    if download: download_pres(pres_id, name)
    utils.extract_images(name)
    pres = Presentation(f'presentations/{name}')
    pres_slides = pres.slides
    session = db_handler.Session()
    current_slide_ids = []
    for i in range(len(pres_slides)):
        with open(f'presentations/images/img_{i}.png', "rb") as image:
            img_bytes = image.read()
            assigned_id = pres_slides[i].notes_slide.notes_text_frame.text
            if assigned_id == '':
                new_slide = Slides(
                    pres_id=pres_id,
                    index=i,
                    thumbnail=img_bytes
                )
                session.add(new_slide)
                session.commit()
                session.refresh(new_slide)
                pres_slides[i].notes_slide.notes_text_frame.text = str(new_slide.id)
                current_slide_ids.append(new_slide.id)
            else:
                slide = session.query(Slides).get(assigned_id)
                slide.index = i
                slide.thumbnail = img_bytes
                current_slide_ids.append(assigned_id)
    pres.save(f'presentations/{name}')
    session.query(Slides).filter(Slides.pres_id == pres_id, ~Slides.id.in_(current_slide_ids)).delete()
    session.commit()
    upload_pres(pres_id, name)


def clear_temp():
    presentations = glob.glob('presentations/*.pptx')
    for pres in presentations:
        os.remove(pres)
    images = glob.glob('presentations/images/*')
    for image in images:
        os.remove(image)


if __name__ == "__main__":
    uvicorn.run("main:app", host="127.0.0.1", port=8000, log_level="info")
