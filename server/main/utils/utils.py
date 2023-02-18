import os.path
import shutil
from pathlib import Path
import cv2
import numpy as np
from aspose import slides
from pptx import Presentation
import aspose.pydrawing as draw

from definitions import ROOT

SROOT = os.path.join(ROOT, 'server/main')

def clear_watermark(path):
    pres = Presentation(path)
    slides = [slide for slide in pres.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    replacements = {
        'Evaluation only.': '',
        'Created with Aspose.Slides for .NET Standard 2.0 22.11.': '',
        'Copyright 2004-2022Aspose Pty Ltd.': ''
    }
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text
    pres.save(path)


def extract_images(application, name):
    pres_path = os.path.join(SROOT, f'presentations/temp/{name}/{name}.pptx')
    img_path = os.path.join(SROOT, f'presentations/temp/{name}/images')
    pres = application.Presentations.Open(pres_path, WithWindow=False)
    slides_num = len(pres.Slides)

    ratios = []
    for i in range(slides_num):
        slide = pres.Slides[i]
        slide_ratio = (slide.CustomLayout.Width / slide.CustomLayout.Height)
        if slide_ratio / (4 / 3) == 0:
            ratios.append(0)
        elif slide_ratio / (16 / 9) == 0:
            ratios.append(1)
        else:
            ratios.append(None)

    pres.Export(img_path, 'png')

    pres.Close()
    pres = None
    return slides_num, ratios


def extract_text(name):
    pres_path = os.path.join(SROOT, f'presentations/temp/{name}/{name}.pptx')
    pres = Presentation(pres_path)
    slides_text = []
    for slide in pres.slides:
        text_buffer = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_buffer += shape.text
        slides_text.append(text_buffer.lower())
    return slides_text


def clear_temp(pres):
    shutil.rmtree(os.path.join(SROOT, f'presentations/temp/{pres.get("id")}'))


def img_eq(db_thumbnail, upd_thumbnail, threshold=0):
    np_arr_db = np.frombuffer(db_thumbnail, np.uint8)
    img_db = cv2.imdecode(np_arr_db, cv2.IMREAD_COLOR)
    np_arr_upd = np.frombuffer(upd_thumbnail, np.uint8)
    img_upd = cv2.imdecode(np.frombuffer(np_arr_upd, np.uint8), cv2.IMREAD_COLOR)

    lab_img_db = cv2.cvtColor(img_db, cv2.COLOR_RGB2Lab)
    lab_img_upd = cv2.cvtColor(img_upd, cv2.COLOR_RGB2Lab)

    h, w, _ = lab_img_db.shape

    if lab_img_db.shape < lab_img_upd.shape:
        h, w, _ = lab_img_db.shape
        lab_img_upd = cv2.resize(lab_img_upd, (w, h))
    elif lab_img_upd.shape < lab_img_db.shape:
        h, w, _ = lab_img_upd.shape
        lab_img_db = cv2.resize(lab_img_db, (w, h))

    img_diff = cv2.subtract(lab_img_upd, lab_img_db)
    mse = np.sum(img_diff ** 2) / (w * h)

    return mse < threshold


def get_pres_thumbnails(pres, n):
    thumbnails = []
    for i in range(n):
        with open(os.path.join(SROOT, f'presentations/temp/{pres.get("id")}/images/Слайд{i + 1}.png'), "rb") as image:
            thumbnails.append(image.read())
    return thumbnails


