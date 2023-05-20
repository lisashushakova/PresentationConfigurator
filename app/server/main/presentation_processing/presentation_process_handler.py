import base64
import os

import pythoncom
import win32com.client
from aspose import slides

from pptx import Presentation

from app.server.main.utils import utils
from app.definitions import SERVER_ROOT
from app.server.main.database.database_handler import Slides


class PresentationProcessHandler:

    def __init__(self, pool):
        self.pool = pool

    def get_slide_ratio(self, slide):
        slide_ratio = (slide.CustomLayout.Width / slide.CustomLayout.Height)
        if abs(slide_ratio - (4 / 3)) < 10e-2:
            return Slides.Ratio.STANDARD_4_TO_3
        elif abs(slide_ratio - (16 / 9)) < 10e-2:
            return Slides.Ratio.WIDESCREEN_16_TO_9
        else:
            raise AttributeError("Unknown slide ratio")

    def crop_presentation(self, presentation, win32com_app_instance):
        pres_path = os.path.join(SERVER_ROOT, f'presentation_processing/temp/{presentation.get("id")}/{presentation.get("id")}.pptx')
        img_path = os.path.join(SERVER_ROOT, f'presentation_processing/temp/{presentation.get("id")}/images')
        os.mkdir(img_path)

        presentation = win32com_app_instance.Presentations.Open(pres_path, WithWindow=False)
        presentation.Export(img_path, 'PNG')

        try:
            presentation_ratio = self.get_slide_ratio(presentation.Slides[0])
            return presentation_ratio
        except IndexError:
            raise AttributeError("Presentation is empty")
        finally:
            presentation.Close()
            presentation = None

    def crop_presentations(self, presentations):
        presentations_ratio = []

        # Start of win32com Application takes significant amount of time
        # so we should previously check if presentation array is not empty
        if len(presentations) > 0:
            pythoncom.CoInitializeEx(0)
            ApplicationPPTX = win32com.client.Dispatch("PowerPoint.Application")

            for result in self.pool.starmap(
                    self.crop_presentation, [(presentation, ApplicationPPTX) for presentation in presentations]):
                presentations_ratio.append(result)

            ApplicationPPTX.Quit()
            ApplicationPPTX = None

        return presentations_ratio

    def __extract_text(self, presentation):
        pres_path = os.path.join(SERVER_ROOT, f'presentation_processing/temp/{presentation.get("id")}/{presentation.get("id")}.pptx')
        pres = Presentation(pres_path)
        text = []
        for slide in pres.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text += shape.text
            text.append(slide_text)
        return text

    def extract_text(self, presentations):
        presentations_text = []
        for result in self.pool.map(self.__extract_text, presentations):
            presentations_text.append(result)
        return presentations_text

    def create_style_template(self, presentation, win32com_app_instance):
        pres_path = os.path.join(
            SERVER_ROOT,
            f"presentation_processing/temp/{presentation.get('id')}/{presentation.get('id')}.pptx"
        )
        style_sample_path = os.path.join(
            SERVER_ROOT,
            f"presentation_processing/styles/style-sample.pptx"
        )
        style_template_directory = os.path.join(
            SERVER_ROOT,
            f"presentation_processing/styles/temp/{presentation.get('id')}"
        )
        style_template_path = os.path.join(style_template_directory, f"{presentation.get('id')}.pptx")

        with slides.Presentation(pres_path) as build_from_pres:
            with slides.Presentation(style_sample_path) as sample_pres:
                # Adding master slide with style
                first_slide_master_slide = build_from_pres.slides[0].layout_slide.master_slide
                new_master_slide = sample_pres.masters.add_clone(first_slide_master_slide)
                sample_pres.slides[0].layout_slide.master_slide = new_master_slide

                # Remove old master
                sample_pres.masters.remove_at(0)

                # Saving template presentation
                if not os.path.exists(style_template_directory):
                    os.mkdir(style_template_directory)
                sample_pres.save(style_template_path, slides.export.SaveFormat.PPTX)
                utils.clear_watermark(style_template_path)

        # Exporting as image
        pres = win32com_app_instance.Presentations.Open(style_template_path, WithWindow=False)
        pres.Export(style_template_directory, 'png')
        pres.Close()
        pres = None

        with open(style_template_directory + '/Слайд1.PNG', "rb") as image:
            style_thumbnail = base64.b64encode(image.read())
            style_template = {
                'id': presentation.get('id'),
                'thumbnail': style_thumbnail
            }

        return style_template

    def create_style_templates(self, presentations):
        # Start of win32com Application takes significant amount of time
        # so we should previously check if presentation array is not empty
        if len(presentations) > 0:
            pythoncom.CoInitializeEx(0)
            ApplicationPPTX = win32com.client.Dispatch("PowerPoint.Application")

            style_path = os.path.join(SERVER_ROOT, "presentation_processing/styles/temp")
            if not os.path.exists(style_path):
                os.mkdir(style_path)

            style_templates = []

            for result in self.pool.starmap(self.create_style_template,
                                            [(presentation, ApplicationPPTX) for presentation in presentations]):
                style_templates.append(result)

            ApplicationPPTX.Quit()
            ApplicationPPTX = None

            return style_templates

        return None


    def build_presentation(self, name, slides_from, ratio, style_template, db_handler, user_id):
        with slides.Presentation() as presentation:
            if ratio == 'widescreen_16_to_9':
                presentation.slide_size.set_size(
                    slides.SlideSizeType.ON_SCREEN_16X9,
                    slides.SlideSizeScaleType.ENSURE_FIT
                )
            else:
                presentation.slide_size.set_size(
                    slides.SlideSizeType.ON_SCREEN,
                    slides.SlideSizeScaleType.ENSURE_FIT
                )

            presentation.slides.remove_at(0)

            seen_thumbs = []
            seen_text = []
            seen_slides = []

            for slide in slides_from:
                with slides.Presentation(
                        os.path.join(SERVER_ROOT, f"presentation_processing/temp/{slide.pres_id}/{slide.pres_id}.pptx")
                ) as pres_from:
                    pres_text = self.__extract_text({'id': slide.pres_id})
                    unique = True
                    slide_thumb = db_handler.get_slide_thumb_by_index(user_id, slide.pres_id, slide.index)
                    slide_text = pres_text[slide.index]
                    for seen_thumb in seen_thumbs:
                        if utils.img_eq(slide_thumb, seen_thumb):
                            unique = False
                            break
                    if unique:
                        presentation.slides.add_clone(pres_from.slides[slide.index])

                        seen_thumbs.append(slide_thumb)
                        seen_text.append(slide_text)
                        seen_slides.append(slide)

            # Applying styles
            if style_template:
                with slides.Presentation(os.path.join(SERVER_ROOT, f"presentation_processing/temp/{style_template}/{style_template}.pptx")) as pres_from:
                    first_slide_master = pres_from.slides[0].layout_slide.master_slide
                    new_master = presentation.masters.add_clone(first_slide_master)
                    for slide in presentation.slides:
                        slide.layout_slide.master_slide = new_master
                    presentation.masters.remove_at(0)

            if not os.path.exists(os.path.join(SERVER_ROOT, f"presentation_processing/built")):
                os.mkdir(os.path.join(SERVER_ROOT, f"presentation_processing/built"))

            if not os.path.exists(os.path.join(SERVER_ROOT, f"presentation_processing/built/{user_id}")):
                os.mkdir(os.path.join(SERVER_ROOT, f"presentation_processing/built/{user_id}"))

            presentation.save(
                os.path.join(SERVER_ROOT, f'presentation_processing/built/{user_id}/{name}.pptx'),
                slides.export.SaveFormat.PPTX
            )
            utils.clear_watermark(os.path.join(SERVER_ROOT, f'presentation_processing/built/{user_id}/{name}.pptx'))

        return seen_text, seen_thumbs, seen_slides