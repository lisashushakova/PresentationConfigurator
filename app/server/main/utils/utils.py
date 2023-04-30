import os
import shutil

import cv2
import numpy as np
from pptx import Presentation

from app.definitions import SERVER_ROOT


def clear_temp():
    temp_path = os.path.join(SERVER_ROOT, f'presentation_processing/temp')
    if os.path.exists(temp_path):
        shutil.rmtree(temp_path)


def clear_styles():
    temp_path = os.path.join(SERVER_ROOT, f'presentation_processing/styles/temp')
    if os.path.exists(temp_path):
        shutil.rmtree(temp_path)


def clear_user_built(user_id):
    temp_path = os.path.join(SERVER_ROOT, f'presentation_processing/built/{user_id}')
    if os.path.exists(temp_path):
        shutil.rmtree(temp_path)


def img_eq(db_thumbnail, upd_thumbnail, crop_bottom=0.1, threshold=0.3):
    np_arr_db = np.frombuffer(db_thumbnail, np.uint8)
    img_db = cv2.imdecode(np_arr_db, cv2.IMREAD_COLOR)
    np_arr_upd = np.frombuffer(upd_thumbnail, np.uint8)
    img_upd = cv2.imdecode(np_arr_upd, cv2.IMREAD_COLOR)

    lab_img_db = cv2.cvtColor(img_db, cv2.COLOR_RGB2Lab)
    lab_img_upd = cv2.cvtColor(img_upd, cv2.COLOR_RGB2Lab)

    h, w, _ = lab_img_db.shape

    if lab_img_db.shape < lab_img_upd.shape:
        h, w, _ = lab_img_db.shape
        lab_img_upd = cv2.resize(lab_img_upd, (w, h))
    elif lab_img_upd.shape < lab_img_db.shape:
        h, w, _ = lab_img_upd.shape
        lab_img_db = cv2.resize(lab_img_db, (w, h))

    img_diff = cv2.subtract(lab_img_upd, lab_img_db)

    # Cropping
    img_diff = img_diff[0:int((1-crop_bottom)*len(img_diff))]

    mse = np.sum(img_diff ** 2) / (w * h)

    return mse < threshold


def clear_watermark(path):
    pres = Presentation(path)
    slides = [slide for slide in pres.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    replacements = {
        'Evaluation only.': '',
        'Created with Aspose.Slides for .NET Standard 2.0 22.11.': '',
        'Copyright 2004-2022Aspose Pty Ltd.': ''
    }
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text
    pres.save(path)
