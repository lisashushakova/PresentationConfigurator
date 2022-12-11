from pathlib import Path
import pythoncom

from pptx import Presentation
import win32com.client


def clear_watermark(path):
    pres = Presentation(path)
    slides = [slide for slide in pres.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    replacements = {
        'Evaluation only.': '',
        'Created with Aspose.Slides for .NET Standard 2.0 22.11.': '',
        'Copyright 2004-2022Aspose Pty Ltd.': ''
    }
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text
    pres.save(path)


def extract_images():
    pythoncom.CoInitializeEx(0)
    Application = win32com.client.Dispatch("PowerPoint.Application")
    pres_path = Path('presentations/temp.pptx').resolve()
    pres = Application.Presentations.Open(pres_path, WithWindow=False)
    slides_num = len(pres.Slides)
    for num in range(slides_num):
        img_path = Path(f'presentations/images/img_{num}.png').resolve()
        pres.Slides[num].Export(img_path, "PNG")
    Application.Quit()
    return slides_num
